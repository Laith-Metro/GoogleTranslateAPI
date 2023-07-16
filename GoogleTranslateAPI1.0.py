import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pygoogletranslation import Translator
from deep_translator import GoogleTranslator
from pptx.util import Inches
import mysql.connector
import tkinter as tk
import random
import os

# Create GUI
def button_click():
    global input_topic, input_count
    input_topic = topic_field.get()
    input_count = int(count_field.get())
    label.config(text="Input retrieved.")
    window.destroy()

window = tk.Tk()
window.resizable(height=100, width=120)
window.geometry('380x200')
window['background'] = 'yellow'
window.title('QuizMaster using Google Translator')
label = tk.Label(window, text="Welcome to QuizMaster", background= "blue")
topic_label = tk.Label(window, text="Please enter a topic:",background= "blue")
count_label = tk.Label(window, text="How many questions should be chosen?" ,background= "blue")
topic_field = tk.Entry(window)
count_field = tk.Entry(window)
button = tk.Button(window, text="Submit",background= "red", command=button_click )

label.pack()
topic_label.pack()
topic_field.pack()
count_label.pack()
count_field.pack()
button.pack()

window.mainloop()

# Connect to db and run query
conn = mysql.connector.connect(
    host='localhost',
    user='root',
    #password='password',
    database='quiz_master'
)

cursor = conn.cursor()
query = "SELECT * FROM questions WHERE topic = '" + input_topic + "'"
cursor.execute(query)

results = cursor.fetchall()
if (len(results) < input_count):
      input_count = len(results)

random_results = random.sample(results, input_count)

# Create a Presentation object
presentation = Presentation()
slide_layout = presentation.slide_layouts[1]
slide_count = 0

for row in random_results:
        slide = presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

# Set the question as the slide title
        title_shape = shapes.title
        question = row[2]
        title_shape.text = question

# Add the potential answers in a textbox and Text Translation using GoogleTranslator API
        answers = [row[3], row[4], row[5], row[6]]


        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        t0 = GoogleTranslator(source='auto', target='ar').translate(answers[0])

        tf.text = "A. " + answers[0] + " <-- Translation --> " + ''.join([str(t0)])

        p = tf.add_paragraph()
        t1=GoogleTranslator(source='auto', target='ar').translate(answers[1])
        p.text = "B. " + answers[1] + " <-- Translation --> " + ''.join([str(t1)])

        p = tf.add_paragraph()
        t2=GoogleTranslator(source='auto', target='ar').translate(answers[2])
        p.text = "C. " + answers[2] + " <-- Translation --> " + ''.join([str(t2)])


        p = tf.add_paragraph()
        t3 = GoogleTranslator(source='auto', target='ar').translate(answers[3])

        p.text = "D. " + answers[3]  +  " <-- Translation --> " + ''.join([str(t3)])
        slide_count += 1

        print(t0)
        print(t1)
        print(t2)
        print(t3)

# Add the code snippet as an image
        img_path = "D:/MyXampp/htdocs/QuizMasterDB/" + row[8]
        if os.path.exists(img_path):
            try:
                left = Inches(1)
                top = Inches(5)
                width = Inches(2)
                height = Inches(2)
                pic = slide.shapes.add_picture(img_path, left, top, width, height)
            except Exception as e:
                 print("Error adding image to slide.")
        else:
             print("Image not found.")

# Save the PowerPoint presentation to a file
presentation.save("C:/Users/Laith/Documents/images/test1.pptx")

cursor.close()
conn.close()